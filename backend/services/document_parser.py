"""
document_parser.py — DocumentParser 防腐层。

降级链：mineru → marker → markitdown → pdfplumber
通过 DOCUMENT_PARSER_BACKEND 配置项控制首选后端。

License 隔离策略：
  - MinerU（AGPL-3.0）：通过 subprocess 调用 CLI，主进程不 import
  - Marker（GPL-3.0）：通过 subprocess 调用 CLI，主进程不 import
  - MarkItDown（MIT）：可直接 import
  - pdfplumber（MIT）：可直接 import，作为最终降级后端
"""
from __future__ import annotations

import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from typing import Protocol

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# 辅助函数
# ---------------------------------------------------------------------------

def _count_headings(markdown: str) -> int:
    """统计 Markdown 中的标题数量（H1/H2/H3）"""
    return len(re.findall(r'^#{1,3}\s', markdown, re.MULTILINE))


def _count_formulas(markdown: str) -> int:
    """估算 Markdown 中的公式数量（$ 符号对数）"""
    dollar_count = len(re.findall(r'\$', markdown))
    return dollar_count // 2


def _is_parse_usable(result: "ParseResult", file_path: str) -> tuple[bool, str]:
    """Guard against parser backends that succeed but return unusable text."""
    text = (result.markdown or "").strip()
    if len(text) < 40:
        return False, "parsed text is too short"

    ext = os.path.splitext(file_path)[1].lower()
    if ext in {".pdf", ".pptx", ".ppt"} and result.page_count >= 3:
        avg_chars = len(text) / max(result.page_count, 1)
        if avg_chars < 20:
            return False, f"parsed text is too sparse ({avg_chars:.1f} chars/page)"

    return True, ""


# ---------------------------------------------------------------------------
# 数据模型
# ---------------------------------------------------------------------------

@dataclass
class ParseResult:
    """文档解析结果。"""
    markdown: str           # 带层级的 Structured Markdown
    heading_count: int      # 检测到的标题数量
    formula_count: int      # 检测到的公式数量
    page_count: int = 0     # 文档页数/幻灯片数
    degraded: bool = False  # 是否降级到低级后端
    backend_used: str = ""  # 实际使用的后端名称


# ---------------------------------------------------------------------------
# 后端协议
# ---------------------------------------------------------------------------

class DocumentParserBackend(Protocol):
    """文档解析后端协议。"""

    def parse(self, file_path: str) -> ParseResult: ...
    def is_available(self) -> bool: ...


# ---------------------------------------------------------------------------
# MinerUBackend（subprocess 隔离，AGPL-3.0）
# ---------------------------------------------------------------------------

class MinerUBackend:
    """
    通过 subprocess 调用 MinerU CLI，主进程不 import mineru。
    规避 AGPL-3.0 传染性。

    调用命令：mineru parse --input <file_path> --output-format json
    stdout JSON 格式：
    {"markdown": "...", "heading_count": 42, "formula_count": 15, "page_count": 100}
    """

    def is_available(self) -> bool:
        """检查 mineru CLI 是否在 PATH 中。"""
        try:
            result = subprocess.run(
                ["mineru", "--version"],
                capture_output=True,
                timeout=5,
            )
            return result.returncode == 0
        except Exception:
            return False

    def parse(self, file_path: str) -> ParseResult:
        """通过 mineru CLI 解析文件，返回 ParseResult。"""
        from config import get_config  # type: ignore[import]
        cfg = get_config()
        timeout = cfg.DOCUMENT_PARSER_TIMEOUT_SECONDS

        try:
            result = subprocess.run(
                ["mineru", "parse", "--input", file_path, "--output-format", "json"],
                capture_output=True,
                text=True,
                timeout=timeout,
            )
        except subprocess.TimeoutExpired as exc:
            raise RuntimeError(
                f"MinerU 解析超时（>{timeout}s）：{file_path}"
            ) from exc
        except Exception as exc:
            raise RuntimeError(f"MinerU 调用失败：{exc}") from exc

        if result.returncode != 0:
            raise RuntimeError(
                f"MinerU 返回非零退出码 {result.returncode}：{result.stderr}"
            )

        try:
            data = json.loads(result.stdout)
        except json.JSONDecodeError as exc:
            raise RuntimeError(
                f"MinerU stdout JSON 解析失败：{exc}\nstdout={result.stdout[:200]}"
            ) from exc

        return ParseResult(
            markdown=data.get("markdown", ""),
            heading_count=int(data.get("heading_count", 0)),
            formula_count=int(data.get("formula_count", 0)),
            page_count=int(data.get("page_count", 0)),
            degraded=False,
            backend_used="mineru",
        )


# ---------------------------------------------------------------------------
# MarkerBackend（subprocess 隔离，GPL-3.0）
# ---------------------------------------------------------------------------

class MarkerBackend:
    """
    通过 subprocess 调用 marker_single CLI，主进程不 import marker。
    规避 GPL-3.0 传染性。

    调用命令：marker_single <file_path> --output_dir <tmpdir> --output_format markdown
    """

    def is_available(self) -> bool:
        """检查 marker_single CLI 是否在 PATH 中。"""
        try:
            result = subprocess.run(
                ["marker_single", "--help"],
                capture_output=True,
                timeout=5,
            )
            return result.returncode == 0
        except Exception:
            return False

    def parse(self, file_path: str) -> ParseResult:
        """通过 marker_single CLI 解析文件，返回 ParseResult。"""
        from config import get_config  # type: ignore[import]
        cfg = get_config()
        timeout = cfg.DOCUMENT_PARSER_TIMEOUT_SECONDS

        tmpdir = tempfile.mkdtemp(prefix="marker_out_")
        try:
            try:
                result = subprocess.run(
                    [
                        "marker_single",
                        file_path,
                        "--output_dir", tmpdir,
                        "--output_format", "markdown",
                    ],
                    capture_output=True,
                    text=True,
                    timeout=timeout,
                )
            except subprocess.TimeoutExpired as exc:
                raise RuntimeError(
                    f"Marker 解析超时（>{timeout}s）：{file_path}"
                ) from exc
            except Exception as exc:
                raise RuntimeError(f"Marker 调用失败：{exc}") from exc

            if result.returncode != 0:
                raise RuntimeError(
                    f"Marker 返回非零退出码 {result.returncode}：{result.stderr}"
                )

            # 查找输出的 .md 文件
            md_content = ""
            for root, _dirs, files in os.walk(tmpdir):
                for fname in files:
                    if fname.endswith(".md"):
                        md_path = os.path.join(root, fname)
                        with open(md_path, "r", encoding="utf-8") as f:
                            md_content = f.read()
                        break
                if md_content:
                    break

            if not md_content:
                raise RuntimeError(
                    f"Marker 未生成 .md 文件，输出目录：{tmpdir}"
                )

            heading_count = _count_headings(md_content)
            formula_count = _count_formulas(md_content)

            return ParseResult(
                markdown=md_content,
                heading_count=heading_count,
                formula_count=formula_count,
                page_count=0,
                degraded=False,
                backend_used="marker",
            )
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)


# ---------------------------------------------------------------------------
# MarkItDownBackend（直接 import，MIT）
# ---------------------------------------------------------------------------

class MarkItDownBackend:
    """
    通过直接 import markitdown 解析文件（MIT 许可）。
    不支持 LaTeX，标记为降级。
    """

    def is_available(self) -> bool:
        try:
            import markitdown  # noqa: F401
            return True
        except ImportError:
            return False

    def parse(self, file_path: str) -> ParseResult:
        from markitdown import MarkItDown  # type: ignore[import]
        md = MarkItDown()
        result = md.convert(file_path)
        markdown = result.text_content
        heading_count = len(re.findall(r'^#{1,3}\s', markdown, re.MULTILINE))
        formula_count = len(re.findall(r'\$', markdown))
        return ParseResult(
            markdown=markdown,
            heading_count=heading_count,
            formula_count=formula_count // 2,  # $ 成对出现
            degraded=True,  # MarkItDown 不支持 LaTeX，标记为降级
            backend_used="markitdown",
        )


# ---------------------------------------------------------------------------
# PdfplumberBackend（直接 import，MIT，最终降级）
# ---------------------------------------------------------------------------

class PdfplumberBackend:
    """
    使用 pdfplumber / python-docx / python-pptx 解析文件（MIT 许可）。
    作为最终降级后端，始终可用。
    """

    def is_available(self) -> bool:
        return True  # 始终可用（pdfplumber 是项目依赖）

    def parse(self, file_path: str) -> ParseResult:
        """根据文件扩展名分发到对应解析器。"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".docx":
            return self._parse_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            return self._parse_pptx(file_path)
        else:
            raise ValueError(f"PdfplumberBackend 不支持的文件格式：{ext}")

    def _parse_pdf(self, file_path: str) -> ParseResult:
        import pdfplumber
        pages: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if not text.strip():
                    logger.warning("第 %d 页无文字内容，跳过", i + 1)
                pages.append(text)
        markdown = "\n\n".join(p for p in pages if p.strip())
        heading_count = _count_headings(markdown)
        formula_count = _count_formulas(markdown)
        return ParseResult(
            markdown=markdown,
            heading_count=heading_count,
            formula_count=formula_count,
            page_count=page_count,
            degraded=True,
            backend_used="pdfplumber",
        )

    def _parse_docx(self, file_path: str) -> ParseResult:
        from docx import Document  # type: ignore[import]
        doc = Document(file_path)
        lines: list[str] = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading 1"):
                lines.append(f"# {para.text}")
            elif style_name.startswith("Heading 2"):
                lines.append(f"## {para.text}")
            elif style_name.startswith("Heading 3"):
                lines.append(f"### {para.text}")
            else:
                lines.append(para.text)
        markdown = "\n\n".join(lines)
        heading_count = _count_headings(markdown)
        formula_count = _count_formulas(markdown)
        return ParseResult(
            markdown=markdown,
            heading_count=heading_count,
            formula_count=formula_count,
            page_count=0,
            degraded=True,
            backend_used="docx",
        )

    def _parse_pptx(self, file_path: str) -> ParseResult:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        sections: list[str] = []
        for n, slide in enumerate(prs.slides, start=1):
            # 提取幻灯片标题
            title_text = ""
            content_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # 判断是否为标题占位符
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
                try:
                    from pptx.util import Pt  # noqa: F401
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    if (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.type in (
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        )
                    ):
                        title_text = shape.text_frame.text.strip()
                        continue
                except Exception:
                    pass
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        content_parts.append(text)

            slide_title = title_text if title_text else f"Slide {n}"
            content = "\n".join(content_parts)
            sections.append(f"## Slide {n}: {slide_title}\n\n{content}")

        markdown = "\n\n".join(sections)
        heading_count = _count_headings(markdown)
        formula_count = _count_formulas(markdown)
        return ParseResult(
            markdown=markdown,
            heading_count=heading_count,
            formula_count=formula_count,
            page_count=slide_count,
            degraded=True,
            backend_used="pptx",
        )


# ---------------------------------------------------------------------------
# DocumentParser 主类
# ---------------------------------------------------------------------------

class DocumentParser:
    """
    文档解析防腐层。

    降级链：mineru → marker → markitdown → pdfplumber
    通过 DOCUMENT_PARSER_BACKEND 配置项控制首选后端。
    """

    # 支持的文件格式
    SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".ppt"}

    def __init__(self, backend: str | None = None) -> None:
        """
        初始化 DocumentParser。

        Args:
            backend: 后端名称（mineru | marker | markitdown | pdfplumber），
                     None 时从 DOCUMENT_PARSER_BACKEND 配置读取。
        """
        if backend is not None:
            self._preferred_backend = backend
        else:
            try:
                from config import get_config  # type: ignore[import]
                cfg = get_config()
                self._preferred_backend = cfg.DOCUMENT_PARSER_BACKEND
            except Exception:
                self._preferred_backend = "mineru"

    def parse(self, file_path: str) -> ParseResult:
        """
        解析文件，返回 ParseResult。

        按降级链依次尝试后端，所有后端失败时使用 PdfplumberBackend 作为最终降级。
        解析完成后记录 INFO 日志。

        Args:
            file_path: 文件路径。

        Returns:
            ParseResult 实例。

        Raises:
            FileNotFoundError: 文件不存在。
            ValueError: 不支持的文件格式。
        """
        # 1. 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")

        # 2. 检查文件格式
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件格式：{ext}")

        filename = os.path.basename(file_path)

        # 3. 按降级链依次尝试后端
        backend_chain = self._get_backend_chain()
        result: ParseResult | None = None

        for backend in backend_chain:
            result = self._try_backend(backend, file_path)
            if result is not None:
                break

        # 4. 所有后端失败时，使用 PdfplumberBackend 作为最终降级
        if result is None:
            logger.error(
                "所有后端均失败，强制使用 PdfplumberBackend 降级解析：%s", filename
            )
            fallback = PdfplumberBackend()
            result = fallback.parse(file_path)

        # 5. 记录 INFO 日志
        logger.info(
            "文档解析完成：文件=%s，格式=%s，页数=%d，标题数=%d，公式数=%d，"
            "后端=%s，降级=%s",
            filename,
            ext,
            result.page_count,
            result.heading_count,
            result.formula_count,
            result.backend_used,
            result.degraded,
        )

        return result

    def _try_backend(
        self,
        backend: DocumentParserBackend,
        file_path: str,
    ) -> ParseResult | None:
        """
        尝试单个后端解析文件。

        Args:
            backend: 后端实例。
            file_path: 文件路径。

        Returns:
            ParseResult 实例，失败时返回 None。
        """
        backend_name = type(backend).__name__
        try:
            if not backend.is_available():
                logger.debug("%s 不可用，跳过", backend_name)
                return None
            result = backend.parse(file_path)
            usable, reason = _is_parse_usable(result, file_path)
            if not usable:
                raise RuntimeError(f"解析结果质量不足：{reason}")
            return result
        except Exception as exc:
            logger.warning(
                "%s 解析失败，降级到下一后端：%s（错误：%s）",
                backend_name,
                os.path.basename(file_path),
                exc,
            )
            return None

    def _get_backend_chain(self) -> list:
        """
        根据 self._preferred_backend 返回有序后端列表。

        首选后端排在最前面，其余按 mineru→marker→markitdown→pdfplumber 顺序。

        Returns:
            后端实例列表。
        """
        all_backends: dict[str, DocumentParserBackend] = {
            "mineru": MinerUBackend(),
            "marker": MarkerBackend(),
            "markitdown": MarkItDownBackend(),
            "pdfplumber": PdfplumberBackend(),
        }

        preferred = self._preferred_backend.lower()
        default_order = ["mineru", "marker", "markitdown", "pdfplumber"]

        # 首选后端排在最前面
        if preferred in all_backends:
            chain_keys = [preferred] + [k for k in default_order if k != preferred]
        else:
            logger.warning(
                "未知的首选后端 '%s'，使用默认降级链", preferred
            )
            chain_keys = default_order

        return [all_backends[k] for k in chain_keys]
