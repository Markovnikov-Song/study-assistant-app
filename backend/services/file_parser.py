"""
file_parser.py — 公共文件解析工具，供 DocumentService 和 ExamService 共用。

PDF/PPTX/DOCX 优先通过 DocumentParser 防腐层解析（结构化 Markdown）；
.txt/.md 保持原有逻辑；降级时记录 WARNING 日志但不中断流程。
"""
from __future__ import annotations

import logging
import os

logger = logging.getLogger(__name__)


def parse_file(tmp_path: str, filename: str) -> str:
    """
    解析文件为文本（Structured Markdown 或纯文本），根据扩展名分发。

    对 PDF/PPTX/DOCX：优先调用 DocumentParser 防腐层，返回带层级的 Markdown。
    对 .txt/.md：直接读取原始文本。
    降级时（ParseResult.degraded=True）记录 WARNING 日志，但不中断流程。

    :param tmp_path: 临时文件路径
    :param filename: 原始文件名（用于判断扩展名）
    :return: 解析后的文本内容
    :raises ValueError: 不支持的文件格式
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext in (".pdf", ".docx", ".pptx", ".ppt"):
        return _parse_via_document_parser(tmp_path, filename)
    elif ext in (".txt", ".md"):
        return _parse_text(tmp_path)
    else:
        raise ValueError(f"不支持的文件格式：{ext}")


def _parse_via_document_parser(tmp_path: str, filename: str) -> str:
    """
    通过 DocumentParser 防腐层解析 PDF/PPTX/DOCX。

    降级时（ParseResult.degraded=True）记录 WARNING 日志。
    """
    try:
        from services.document_parser import DocumentParser
        parser = DocumentParser()
        result = parser.parse(tmp_path)
        if result.degraded:
            logger.warning(
                "文件 %s 解析降级（后端：%s），结构化信息可能不完整",
                filename,
                result.backend_used,
            )
        return result.markdown
    except Exception as e:
        # DocumentParser 完全失败时，回退到原有逻辑
        logger.error(
            "DocumentParser 解析失败（%s），回退到原生解析器：%s", filename, e
        )
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return _parse_pdf_fallback(tmp_path)
        elif ext == ".docx":
            return _parse_docx_fallback(tmp_path)
        elif ext in (".pptx", ".ppt"):
            return _parse_pptx_fallback(tmp_path)
        raise


def _parse_pdf_fallback(tmp_path: str) -> str:
    """pdfplumber 纯文本降级（DocumentParser 完全失败时使用）。"""
    import pdfplumber
    pages: list[str] = []
    with pdfplumber.open(tmp_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if not text.strip():
                logger.warning("第 %d 页无文字内容，跳过", i)
            pages.append(text)
    return "\n".join(pages)


def _parse_docx_fallback(tmp_path: str) -> str:
    """python-docx 纯文本降级（DocumentParser 完全失败时使用）。"""
    from docx import Document
    doc = Document(tmp_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _parse_pptx_fallback(tmp_path: str) -> str:
    """python-pptx 纯文本降级（DocumentParser 完全失败时使用）。"""
    from pptx import Presentation
    prs = Presentation(tmp_path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
    return "\n".join(texts)


def _parse_text(tmp_path: str) -> str:
    with open(tmp_path, "r", encoding="utf-8") as f:
        return f.read()
